from fastapi import FastAPI, UploadFile, File, Depends, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from pydantic import BaseModel
from openai import OpenAI
from dotenv import load_dotenv
from sqlalchemy.orm import Session
from jose import JWTError, jwt
from datetime import timedelta
import os, sys, tempfile

sys.path.append(os.path.dirname(__file__))
from pdf_extractor import extract_text_from_pdf
from database import get_db, User, ChatHistory, SavedContent
from auth import (authenticate_user, create_user, create_access_token,
                  get_user_by_username, get_user_by_email, SECRET_KEY, ALGORITHM,
                  ACCESS_TOKEN_EXPIRE_MINUTES)

load_dotenv()

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

client = OpenAI(
    api_key=os.getenv("GROQ_API_KEY"),
    base_url="https://api.groq.com/openai/v1"
)

stored_text = {}
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="token", auto_error=False)

LANGUAGE_INSTRUCTION = """
Detect the language of the user's message and respond in that exact same language.
Supported: Nepali, Hindi, English, Spanish, French, German, Italian, Portuguese,
Russian, Chinese, Japanese, Korean, Arabic, Turkish, Dutch, Polish, Swedish,
Norwegian, Danish, Finnish, Bengali, Urdu, Tagalog, Cebuano, Visayan,
Malay, Indonesian, Vietnamese, Thai, Swahili.
Never mix languages. Always match the user's language exactly.
"""

# ─── Auth helpers ──────────────────────────────────────────────────────────────

def get_current_user(token: str = Depends(oauth2_scheme), db: Session = Depends(get_db)):
    if not token:
        return None
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        username = payload.get("sub")
        if not username:
            return None
        return get_user_by_username(db, username)
    except JWTError:
        return None

# ─── Pydantic models ───────────────────────────────────────────────────────────

class SignupRequest(BaseModel):
    username: str
    email: str
    password: str

class ChatRequest(BaseModel):
    session_id: str
    question: str

class ActionRequest(BaseModel):
    session_id: str
    count: int = 10

class EssayRequest(BaseModel):
    session_id: str
    essay_text: str

class HomeworkRequest(BaseModel):
    session_id: str
    question: str
    subject: str = "general"

class Eli5Request(BaseModel):
    session_id: str
    topic: str

class DebateRequest(BaseModel):
    session_id: str
    topic: str

class CompareRequest(BaseModel):
    session_id_1: str
    session_id_2: str

class SaveRequest(BaseModel):
    type: str
    title: str
    content: str
    session_id: str

# ─── Auth routes ───────────────────────────────────────────────────────────────

@app.post("/signup")
def signup(req: SignupRequest, db: Session = Depends(get_db)):
    try:
        if get_user_by_username(db, req.username):
            raise HTTPException(status_code=400, detail="Username already taken")
        if get_user_by_email(db, req.email):
            raise HTTPException(status_code=400, detail="Email already registered")
        user = create_user(db, req.username, req.email, req.password)
        token = create_access_token(
            {"sub": user.username},
            timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
        )
        return {"access_token": token, "token_type": "bearer", "username": user.username}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Signup failed: {str(e)}")

@app.post("/token")
def login(form: OAuth2PasswordRequestForm = Depends(), db: Session = Depends(get_db)):
    try:
        user = authenticate_user(db, form.username, form.password)
        if not user:
            raise HTTPException(status_code=401, detail="Wrong username or password")
        token = create_access_token(
            {"sub": user.username},
            timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
        )
        return {"access_token": token, "token_type": "bearer", "username": user.username}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Login failed: {str(e)}")

@app.get("/me")
def me(current_user: User = Depends(get_current_user)):
    if not current_user:
        raise HTTPException(status_code=401, detail="Not logged in")
    return {"id": current_user.id, "username": current_user.username, "email": current_user.email}

# ─── Save routes ───────────────────────────────────────────────────────────────

@app.post("/save")
def save_content(req: SaveRequest, current_user: User = Depends(get_current_user), db: Session = Depends(get_db)):
    if not current_user:
        raise HTTPException(status_code=401, detail="Login required")
    item = SavedContent(
        user_id=current_user.id,
        type=req.type,
        title=req.title,
        content=req.content,
        session_id=req.session_id
    )
    db.add(item)
    db.commit()
    return {"message": "Saved successfully"}

@app.get("/saved")
def get_saved(current_user: User = Depends(get_current_user), db: Session = Depends(get_db)):
    if not current_user:
        raise HTTPException(status_code=401, detail="Login required")
    items = db.query(SavedContent).filter(
        SavedContent.user_id == current_user.id
    ).order_by(SavedContent.created_at.desc()).all()
    return [{"id": i.id, "type": i.type, "title": i.title, "content": i.content, "session_id": i.session_id} for i in items]

@app.delete("/saved/{item_id}")
def delete_saved(item_id: int, current_user: User = Depends(get_current_user), db: Session = Depends(get_db)):
    if not current_user:
        raise HTTPException(status_code=401, detail="Login required")
    item = db.query(SavedContent).filter(
        SavedContent.id == item_id,
        SavedContent.user_id == current_user.id
    ).first()
    if not item:
        raise HTTPException(status_code=404, detail="Not found")
    db.delete(item)
    db.commit()
    return {"message": "Deleted"}

# ─── File upload ───────────────────────────────────────────────────────────────

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    contents = await file.read()
    filename = file.filename.lower()
    text = ""

    try:
        if filename.endswith(".pdf"):
            text = extract_text_from_pdf(contents)

        elif filename.endswith(".docx"):
            import docx2txt
            with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as f:
                f.write(contents); tmp = f.name
            text = docx2txt.process(tmp)
            os.unlink(tmp)

        elif filename.endswith(".pptx"):
            from pptx import Presentation
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as f:
                f.write(contents); tmp = f.name
            prs = Presentation(tmp)
            text = "\n".join([
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            ])
            os.unlink(tmp)

        elif filename.endswith(".csv"):
            text = contents.decode("utf-8")

        elif filename.endswith((".png", ".jpg", ".jpeg", ".webp")):
            import base64
            ext  = filename.split(".")[-1]
            mime = "image/jpeg" if ext in ["jpg", "jpeg"] else f"image/{ext}"
            b64  = base64.b64encode(contents).decode()
            response = client.chat.completions.create(
                model="llama-3.2-11b-vision-preview",
                messages=[{"role": "user", "content": [
                    {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
                    {"type": "text", "text": "Extract and transcribe all text from this image accurately."}
                ]}],
                max_tokens=2000
            )
            text = response.choices[0].message.content

        else:
            text = contents.decode("utf-8", errors="ignore")

    except Exception as e:
        text = contents.decode("utf-8", errors="ignore")

    if not text or len(text.strip()) < 10:
        return {"error": "Could not extract text from file", "session_id": None}

    session_id = file.filename.replace(" ", "_")
    stored_text[session_id] = text
    return {
        "session_id": session_id,
        "message": "Uploaded successfully",
        "preview": text[:300],
        "word_count": len(text.split())
    }

# ─── Helper ────────────────────────────────────────────────────────────────────

def ask(system: str, user: str, max_tokens: int = 4000) -> str:
    response = client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[
            {"role": "system", "content": system},
            {"role": "user",   "content": user}
        ],
        max_tokens=max_tokens
    )
    return response.choices[0].message.content

def get_notes(session_id: str) -> str:
    return stored_text.get(session_id, "")

# ─── AI endpoints ──────────────────────────────────────────────────────────────

@app.post("/chat")
async def chat(req: ChatRequest):
    notes = get_notes(req.session_id)
    if not notes:
        return {"answer": "No notes found. Please upload a file first."}
    return {"answer": ask(
        f"{LANGUAGE_INSTRUCTION}\nYou are a helpful study assistant. Answer questions based on these notes:\n\n{notes}\n\nDetect the user's language and respond in that same language.",
        req.question
    )}

@app.post("/summarize")
async def summarize(req: ActionRequest):
    notes = get_notes(req.session_id)
    if not notes: return {"summary": "No notes found."}
    return {"summary": ask(
        f"{LANGUAGE_INSTRUCTION}\nSummarize these notes with key points and main ideas. Use the same language as the notes.",
        notes
    )}

@app.post("/quiz")
async def quiz(req: ActionRequest):
    notes = get_notes(req.session_id)
    if not notes: return {"quiz": "No notes found."}
    return {"quiz": ask(
        f"{LANGUAGE_INSTRUCTION}\nGenerate exactly {req.count} multiple choice questions.\nFormat:\nQ: question\nA) option\nB) option\nC) option\nD) option\nAnswer: X\n\nSeparate with blank lines. Use same language as notes.",
        notes
    )}

@app.post("/flashcards")
async def flashcards(req: ActionRequest):
    notes = get_notes(req.session_id)
    if not notes: return {"flashcards": "No notes found."}
    return {"flashcards": ask(
        f"{LANGUAGE_INSTRUCTION}\nCreate exactly {req.count} flashcards.\nFormat:\nFront: question\nBack: answer\n---\nUse same language as notes.",
        notes
    )}

@app.post("/exam_predictor")
async def exam_predictor(req: ActionRequest):
    notes = get_notes(req.session_id)
    if not notes: return {"exam_predictor": "No notes found."}
    return {"exam_predictor": ask(
        f"{LANGUAGE_INSTRUCTION}\nYou are an expert exam predictor. Based on these notes, predict the {req.count} most likely exam questions. For each question explain WHY it's likely to appear. Use same language as notes.",
        notes
    )}

@app.post("/study_plan")
async def study_plan(req: ActionRequest):
    notes = get_notes(req.session_id)
    if not notes: return {"study_plan": "No notes found."}
    return {"study_plan": ask(
        f"{LANGUAGE_INSTRUCTION}\nCreate a detailed 7-day study plan based on these notes. For each day specify: topics to cover, activities, time needed, and goals. Use same language as notes.",
        notes
    )}

@app.post("/key_terms")
async def key_terms(req: ActionRequest):
    notes = get_notes(req.session_id)
    if not notes: return {"key_terms": "No notes found."}
    return {"key_terms": ask(
        f"{LANGUAGE_INSTRUCTION}\nExtract the {req.count} most important key terms and concepts. For each: the term, definition, and why it's important. Use same language as notes.",
        notes
    )}

@app.post("/mind_map")
async def mind_map(req: ActionRequest):
    notes = get_notes(req.session_id)
    if not notes: return {"mind_map": "No notes found."}
    return {"mind_map": ask(
        f"{LANGUAGE_INSTRUCTION}\nCreate a detailed mind map outline.\nFormat:\nMain Topic\n  - Subtopic\n    * Detail\n    * Detail\n  - Subtopic\n    * Detail\nUse same language as notes.",
        notes
    )}

@app.post("/eli5")
async def eli5(req: Eli5Request):
    notes = get_notes(req.session_id)
    topic = req.topic or "the main concept"
    context = f"Using these notes as context:\n{notes}\n\n" if notes else ""
    return {"eli5": ask(
        f"{LANGUAGE_INSTRUCTION}\nExplain '{topic}' as simply as possible, like explaining to a 5 year old. Use simple words, fun analogies, and real life examples. Detect language from the topic.",
        f"{context}Explain: {topic}"
    )}

@app.post("/formula_sheet")
async def formula_sheet(req: ActionRequest):
    notes = get_notes(req.session_id)
    if not notes: return {"formula_sheet": "No notes found."}
    return {"formula_sheet": ask(
        f"{LANGUAGE_INSTRUCTION}\nExtract all formulas, equations, and mathematical expressions. Format as a clean reference sheet with: formula name, the formula, variables explained, and when to use it. Use same language as notes.",
        notes
    )}

@app.post("/chapter_summary")
async def chapter_summary(req: ActionRequest):
    notes = get_notes(req.session_id)
    if not notes: return {"chapter_summary": "No notes found."}
    return {"chapter_summary": ask(
        f"{LANGUAGE_INSTRUCTION}\nDivide these notes into chapters or sections and summarize each one separately with clear headings. Use same language as notes.",
        notes
    )}

@app.post("/simplify_words")
async def simplify_words(req: ActionRequest):
    notes = get_notes(req.session_id)
    if not notes: return {"simplify_words": "No notes found."}
    return {"simplify_words": ask(
        f"{LANGUAGE_INSTRUCTION}\nRewrite these notes in the simplest possible language. Replace difficult words with easy ones. Keep all the information but make it easy to understand. Use same language as notes.",
        notes
    )}

@app.post("/fill_blanks")
async def fill_blanks(req: ActionRequest):
    notes = get_notes(req.session_id)
    if not notes: return {"fill_blanks": "No notes found."}
    return {"fill_blanks": ask(
        f"{LANGUAGE_INSTRUCTION}\nCreate exactly {req.count} fill-in-the-blank exercises.\nFormat:\n1. Sentence with _______ blank\nAnswer: missing word\n\nUse same language as notes.",
        notes
    )}

@app.post("/true_false")
async def true_false(req: ActionRequest):
    notes = get_notes(req.session_id)
    if not notes: return {"true_false": "No notes found."}
    return {"true_false": ask(
        f"{LANGUAGE_INSTRUCTION}\nCreate exactly {req.count} true/false questions.\nFormat:\n1. Statement\nAnswer: True/False\nExplanation: why\n\nMix true and false equally. Use same language as notes.",
        notes
    )}

@app.post("/short_answer")
async def short_answer(req: ActionRequest):
    notes = get_notes(req.session_id)
    if not notes: return {"short_answer": "No notes found."}
    return {"short_answer": ask(
        f"{LANGUAGE_INSTRUCTION}\nCreate exactly {req.count} short answer questions.\nFormat:\nQ: question\nA: answer (2-3 sentences)\n\nUse same language as notes.",
        notes
    )}

@app.post("/debate")
async def debate(req: DebateRequest):
    notes = get_notes(req.session_id)
    topic = req.topic or "the main topic from the notes"
    context = f"Using these notes as context:\n{notes}\n\n" if notes else ""
    return {"debate": ask(
        f"{LANGUAGE_INSTRUCTION}\nPresent a balanced debate on '{topic}'.\nFormat:\n## Arguments FOR\n- point\n\n## Arguments AGAINST\n- point\n\n## Conclusion\nBalanced summary\n\nDetect language from topic.",
        f"{context}Debate topic: {topic}"
    )}

@app.post("/essay_grade")
async def essay_grade(req: EssayRequest):
    notes = get_notes(req.session_id)
    context = f"Reference notes:\n{notes}\n\n" if notes else ""
    return {"essay_grade": ask(
        f"{LANGUAGE_INSTRUCTION}\nYou are an expert essay grader. Grade this essay and provide detailed feedback.\nFormat:\n## Overall Grade: X/10\n\n## Strengths\n- point\n\n## Areas to Improve\n- point\n\n## Detailed Feedback\nParagraph by paragraph analysis\n\n## Suggestions\nSpecific improvements\n\nDetect language from essay.",
        f"{context}Essay to grade:\n{req.essay_text}"
    )}

@app.post("/homework_help")
async def homework_help(req: HomeworkRequest):
    notes = get_notes(req.session_id)
    context = f"Reference notes:\n{notes}\n\n" if notes else ""
    return {"homework_help": ask(
        f"{LANGUAGE_INSTRUCTION}\nYou are an expert {req.subject} tutor. Help solve this problem step by step.\nFormat:\n## Understanding the Problem\n\n## Step-by-Step Solution\nStep 1: ...\nStep 2: ...\n\n## Answer\n\n## Key Concepts Used\n\nDetect language from question.",
        f"{context}Subject: {req.subject}\nQuestion: {req.question}"
    )}

@app.post("/compare")
async def compare(req: CompareRequest):
    notes1 = get_notes(req.session_id_1)
    notes2 = get_notes(req.session_id_2)
    if not notes1 or not notes2:
        return {"compare": "Please upload two documents to compare."}
    return {"compare": ask(
        f"{LANGUAGE_INSTRUCTION}\nCompare these two documents thoroughly.\nFormat:\n## Document 1 Summary\n\n## Document 2 Summary\n\n## Key Similarities\n- point\n\n## Key Differences\n- point\n\n## Which is More Comprehensive?\n\n## Conclusion\n\nUse English for comparison.",
        f"Document 1:\n{notes1[:3000]}\n\nDocument 2:\n{notes2[:3000]}"
    )}

# ─── Health check ──────────────────────────────────────────────────────────────

@app.get("/")
def root():
    return {"message": "AI Student Assistant API is running", "status": "ok"}